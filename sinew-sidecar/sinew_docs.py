"""
Sinew document sidecar — surgical read/edit for DOCX, PDF, XLSX, PPTX.
Protocol: reads one JSON object from stdin, writes one JSON object to stdout, then exits.
"""
import sys
import json
import os
from pathlib import Path


def main():
    raw = sys.stdin.read().strip()
    if not raw:
        _out({"ok": False, "error": "empty input"})
        return
    try:
        cmd = json.loads(raw)
    except json.JSONDecodeError as e:
        _out({"ok": False, "error": f"invalid JSON: {e}"})
        return

    try:
        result = dispatch(cmd)
    except ImportError as e:
        pkg = str(e).split("'")[1] if "'" in str(e) else str(e)
        result = {
            "ok": False,
            "error": (
                f"Missing Python package: {pkg}. "
                f"Run: pip install python-docx PyMuPDF openpyxl python-pptx"
            ),
        }
    except Exception as e:
        result = {"ok": False, "error": str(e)}

    _out(result)


def _out(obj):
    print(json.dumps(obj, ensure_ascii=False), flush=True)


def dispatch(cmd):
    op = cmd.get("op")
    path = cmd.get("path", "")
    if not path:
        return {"ok": False, "error": "missing path"}

    ext = Path(path).suffix.lower()

    if op == "read":
        return _read(path, ext)
    elif op == "find_replace":
        find = cmd.get("find")
        replace = cmd.get("replace", "")
        if find is None:
            return {"ok": False, "error": "missing 'find'"}
        return _find_replace(path, ext, find, replace, cmd.get("all", True))
    elif op == "insert_after":
        return _insert_after(path, ext, cmd.get("after", ""), cmd.get("content", ""))
    elif op == "delete_paragraph":
        return _delete_paragraph(path, ext, cmd.get("paragraph", ""))
    else:
        return {"ok": False, "error": f"unknown op: {op}"}


# ---------------------------------------------------------------------------
# READ
# ---------------------------------------------------------------------------

def _read(path, ext):
    if ext == ".docx":
        return _read_docx(path)
    elif ext == ".pdf":
        return _read_pdf(path)
    elif ext in (".xlsx", ".xls"):
        return _read_xlsx(path)
    elif ext in (".pptx", ".ppt"):
        return _read_pptx(path)
    else:
        # Plain text fallback
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            return {"ok": True, "content": f.read()}


def _read_docx(path):
    import docx
    doc = docx.Document(path)
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells)
            if row_text.strip():
                parts.append(row_text)
    return {"ok": True, "content": "\n".join(parts)}


def _read_pdf(path):
    import fitz
    doc = fitz.open(path)
    parts = []
    for i, page in enumerate(doc):
        text = page.get_text().strip()
        if text:
            parts.append(f"--- Page {i + 1} ---\n{text}")
    return {"ok": True, "content": "\n\n".join(parts)}


def _read_xlsx(path):
    import openpyxl
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    parts = []
    for name in wb.sheetnames:
        ws = wb[name]
        parts.append(f"## {name}")
        for row in ws.iter_rows(values_only=True):
            if any(c is not None for c in row):
                parts.append("\t".join("" if c is None else str(c) for c in row))
    return {"ok": True, "content": "\n".join(parts)}


def _read_pptx(path):
    from pptx import Presentation
    prs = Presentation(path)
    parts = []
    for i, slide in enumerate(prs.slides):
        parts.append(f"## Slide {i + 1}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        parts.append(t)
    return {"ok": True, "content": "\n".join(parts)}


# ---------------------------------------------------------------------------
# FIND / REPLACE
# ---------------------------------------------------------------------------

def _find_replace(path, ext, find, replace, replace_all):
    if ext == ".docx":
        return _fr_docx(path, find, replace, replace_all)
    elif ext == ".pdf":
        return _fr_pdf(path, find, replace, replace_all)
    elif ext in (".xlsx", ".xls"):
        return _fr_xlsx(path, find, replace, replace_all)
    elif ext in (".pptx", ".ppt"):
        return _fr_pptx(path, find, replace, replace_all)
    else:
        return {"ok": False, "error": f"find_replace not supported for {ext}"}


def _fr_docx(path, find, replace, replace_all):
    import docx
    doc = docx.Document(path)
    count = 0

    def _replace_para(para):
        nonlocal count
        if find not in para.text:
            return
        # Rebuild runs preserving formatting of the first run that contains the match
        full = "".join(r.text for r in para.runs)
        if find not in full:
            return
        n = full.count(find) if replace_all else 1
        new_text = full.replace(find, replace) if replace_all else full.replace(find, replace, 1)
        # Put all text in run[0], clear others — preserves run[0] formatting
        if para.runs:
            para.runs[0].text = new_text
            for r in para.runs[1:]:
                r.text = ""
        count += n

    for para in doc.paragraphs:
        _replace_para(para)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    _replace_para(para)

    doc.save(path)
    return {"ok": True, "replacements": count}


def _fr_pdf(path, find, replace, replace_all):
    import fitz
    doc = fitz.open(path)
    count = 0
    for page in doc:
        instances = page.search_for(find)
        if not instances:
            continue
        for rect in instances:
            page.add_redact_annot(rect, replace)
            count += 1
            if not replace_all:
                break
        page.apply_redactions()
        if not replace_all and count > 0:
            break
    # Save in-place with incremental update to preserve max structure
    doc.save(path, incremental=True, encryption=fitz.PDF_ENCRYPT_KEEP)
    return {"ok": True, "replacements": count}


def _fr_xlsx(path, find, replace, replace_all):
    import openpyxl
    wb = openpyxl.load_workbook(path)
    count = 0
    done = False
    for name in wb.sheetnames:
        if done:
            break
        ws = wb[name]
        for row in ws.iter_rows():
            if done:
                break
            for cell in row:
                if isinstance(cell.value, str) and find in cell.value:
                    if replace_all:
                        c = cell.value.count(find)
                        cell.value = cell.value.replace(find, replace)
                        count += c
                    else:
                        cell.value = cell.value.replace(find, replace, 1)
                        count += 1
                        done = True
                        break
    wb.save(path)
    return {"ok": True, "replacements": count}


def _fr_pptx(path, find, replace, replace_all):
    from pptx import Presentation
    prs = Presentation(path)
    count = 0
    done = False
    for slide in prs.slides:
        if done:
            break
        for shape in slide.shapes:
            if done:
                break
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if done:
                        break
                    for run in para.runs:
                        if find in run.text:
                            if replace_all:
                                c = run.text.count(find)
                                run.text = run.text.replace(find, replace)
                                count += c
                            else:
                                run.text = run.text.replace(find, replace, 1)
                                count += 1
                                done = True
                                break
    prs.save(path)
    return {"ok": True, "replacements": count}


# ---------------------------------------------------------------------------
# INSERT AFTER
# ---------------------------------------------------------------------------

def _insert_after(path, ext, after_text, content):
    if ext == ".docx":
        return _insert_after_docx(path, after_text, content)
    else:
        return {"ok": False, "error": f"insert_after not supported for {ext}"}


def _insert_after_docx(path, after_text, content):
    import docx
    from docx.oxml.ns import qn
    from copy import deepcopy
    import lxml.etree as etree

    doc = docx.Document(path)
    target = None
    for para in doc.paragraphs:
        if after_text in para.text:
            target = para
            break
    if target is None:
        return {"ok": False, "error": f"Paragraph containing '{after_text}' not found"}

    # Insert new paragraph after target using lxml
    new_para = docx.oxml.OxmlElement("w:p")
    new_r = docx.oxml.OxmlElement("w:r")
    new_t = docx.oxml.OxmlElement("w:t")
    new_t.text = content
    new_r.append(new_t)
    new_para.append(new_r)
    target._element.addnext(new_para)

    doc.save(path)
    return {"ok": True}


# ---------------------------------------------------------------------------
# DELETE PARAGRAPH
# ---------------------------------------------------------------------------

def _delete_paragraph(path, ext, paragraph_text):
    if ext == ".docx":
        return _delete_paragraph_docx(path, paragraph_text)
    else:
        return {"ok": False, "error": f"delete_paragraph not supported for {ext}"}


def _delete_paragraph_docx(path, paragraph_text):
    import docx
    doc = docx.Document(path)
    found = False
    for para in doc.paragraphs:
        if paragraph_text in para.text:
            p = para._element
            p.getparent().remove(p)
            found = True
            break
    if not found:
        return {"ok": False, "error": f"Paragraph containing '{paragraph_text}' not found"}
    doc.save(path)
    return {"ok": True}


if __name__ == "__main__":
    main()
